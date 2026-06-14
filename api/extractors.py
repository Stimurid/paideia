"""Text extractors для PDF/DOCX/PPTX/MD/TXT/URL.

Возвращают чистый текст (с сохранёнными разделителями страниц/слайдов) +
метаданные источника. Парсинг локальный, без LLM.
"""

from __future__ import annotations

import io
import re
from pathlib import Path
from typing import Any

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

try:
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    _EPUB_OK = True
except ImportError:
    _EPUB_OK = False


def extract_pdf(data: bytes) -> dict[str, Any]:
    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                pages.append(f"\n\n--- page {i} ---\n\n{t}")
    return {"format": "pdf", "page_count": len(pages),
            "text": "\n".join(pages).strip()}


def extract_docx(data: bytes) -> dict[str, Any]:
    doc = DocxDocument(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # таблицы — отдельным проходом, плоско
    table_rows: list[str] = []
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
    text = "\n".join(paragraphs)
    if table_rows:
        text += "\n\n--- tables ---\n\n" + "\n".join(table_rows)
    return {"format": "docx", "paragraph_count": len(paragraphs),
            "table_row_count": len(table_rows), "text": text.strip()}


def extract_pptx(data: bytes) -> dict[str, Any]:
    prs = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                if t.strip():
                    chunks.append(t.strip())
            if shape.has_table:
                tbl_rows = []
                for row in shape.table.rows:
                    tbl_rows.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
                if tbl_rows:
                    chunks.append("table:\n" + "\n".join(tbl_rows))
        if chunks:
            slides.append(f"\n\n--- slide {i} ---\n\n" + "\n\n".join(chunks))
    return {"format": "pptx", "slide_count": len(slides),
            "text": "\n".join(slides).strip()}


def extract_md_or_txt(data: bytes) -> dict[str, Any]:
    text = data.decode("utf-8", errors="replace").strip()
    return {"format": "md", "text": text}


def extract_epub(data: bytes) -> dict[str, Any]:
    """EPUB → plain text. Каждая глава (item) разделена маркером."""
    if not _EPUB_OK:
        raise RuntimeError("EbookLib/bs4 not installed — добавь в requirements.txt")
    # ebooklib читает только с диска; пишем во временный файл
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".epub", delete=False) as tf:
        tf.write(data)
        tmp_path = tf.name
    try:
        book = epub.read_epub(tmp_path)
    finally:
        Path(tmp_path).unlink(missing_ok=True)

    chapters: list[str] = []
    for i, item in enumerate(book.get_items_of_type(ebooklib.ITEM_DOCUMENT), 1):
        try:
            html = item.get_content()
            soup = BeautifulSoup(html, "html.parser")
            for tag in soup(["script", "style"]):
                tag.decompose()
            text = soup.get_text(separator="\n").strip()
            text = "\n".join(line.strip() for line in text.split("\n") if line.strip())
            if text:
                chapters.append(f"\n\n--- chapter {i} ---\n\n{text}")
        except Exception:
            continue

    return {"format": "epub", "chapter_count": len(chapters),
            "text": "\n".join(chapters).strip()}


_EXTENSIONS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".md": extract_md_or_txt,
    ".txt": extract_md_or_txt,
    ".epub": extract_epub,
}


def extract_file(filename: str, data: bytes) -> dict[str, Any]:
    ext = Path(filename).suffix.lower()
    fn = _EXTENSIONS.get(ext)
    if not fn:
        raise ValueError(f"unsupported extension: {ext}")
    out = fn(data)
    out["filename"] = filename
    out["bytes"] = len(data)
    return out


# ---------------------------------------------------------------------------
# Сегментация длинного текста на «вероятные кейсы»
# ---------------------------------------------------------------------------


# Эвристика: в нашем корпусе кейсы обычно начинаются с заголовка третьего уровня
# (## или ###), таблицы строк со столбцом «организация», или маркера слайда с
# названием. Для длинных документов делаем грубое разрезание на чанки ~2000-
# 4000 символов с overlap, потом LLM решает, кейс ли каждый сегмент.

_CASE_HEADERS = re.compile(
    r"(?m)^(?:#{2,4}\s+|\*\*\s*[А-ЯA-Z][^*]{5,120}\*\*\s*$|---\s+(?:page|slide)\s+\d+\s+---)"
)


def segment_into_candidates(text: str, *, max_chunk: int = 3500,
                            overlap: int = 200) -> list[dict[str, Any]]:
    """Разбить длинный текст на сегменты-кандидаты в кейсы.

    Стратегия:
    1. Найти все эвристические заголовки.
    2. Если их много (>= 5) — режем по заголовкам.
    3. Если мало — режем на чанки фиксированного размера.
    """
    text = text.strip()
    if not text:
        return []

    headers = list(_CASE_HEADERS.finditer(text))
    segments: list[dict[str, Any]] = []

    if len(headers) >= 5:
        positions = [m.start() for m in headers] + [len(text)]
        for i in range(len(positions) - 1):
            start, end = positions[i], positions[i + 1]
            chunk = text[start:end].strip()
            if len(chunk) > 100:
                segments.append({
                    "kind": "header_section",
                    "header": text[start:start + 120].split("\n")[0].strip(),
                    "text": chunk,
                    "offset": start,
                })
    else:
        step = max_chunk - overlap
        for start in range(0, len(text), step):
            chunk = text[start:start + max_chunk]
            if len(chunk) > 100:
                segments.append({
                    "kind": "fixed_chunk",
                    "header": chunk[:80].replace("\n", " ").strip() + "…",
                    "text": chunk,
                    "offset": start,
                })

    return segments
